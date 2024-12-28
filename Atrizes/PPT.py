from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
# from pptx.enum.text import MSO_ANCHOR
#from pptx.util import PP_ALIGN  # Import this for text alignment
from pptx.enum.text import PP_ALIGN
from re import sub
import sys

# Insert the path of modules folder  
sys.path.insert(0, "D:\\iTunes\\Novelas") 
sys.path.insert(0, "D:\\iTunes\\Codes")
import Excel # type: ignore
# from Images import Wid # type: ignore

import pandas as pd

# Disable the SettingWithCopyWarning
pd.options.mode.chained_assignment = None

# Create a new PowerPoint presentation
presentation = Presentation()

# RETURNS THE LIST OF NOVELAS FROM A GIVEN ACTOR
def read_novelas(df, atriz):
    # Filter for Atriz 'X'
    filt_df = df[df['Atriz'] == atriz]

    # Create the new column by concatenating 'Ano' and 'Novela'
    filt_df['Ano-Novela'] = filt_df['Ano'].astype(str) + " " + filt_df['Novela']

    # Sort by 'Ano'
    sorted_df = filt_df.sort_values(by='Ano')

    # Create the final list
    res_list = sorted_df['Ano-Novela'].tolist()
    res_list = [sub(r"\s+", " ", x.strip()) for x in res_list]

    return res_list

def add_slide(atriz, img_path, novelas):
    # Set slide dimensions for a 16:9 aspect ratio (13.33 inches by 7.5 inches)
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5.625)

    # Add a blank slide layout to customize
    slide_layout = presentation.slide_layouts[5]  # Using a blank layout
    slide = presentation.slides.add_slide(slide_layout)

    # Add title at the top
    title = slide.shapes.title
    title.text = atriz

    # Set position and size of the title box
    title.left = Inches(0.3)  # Horizontal position from the left
    title.top = Inches(0.05)    # Vertical position from the top (era 0.1)
    title.width = Inches(6.33)  # Set width
    title.height = Inches(0.64)  # Set height

    # Rotate the title box
    title.rotation = 0  # No rotation

    title.text_frame.paragraphs[0].font.name = "Bookman Old Style"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Set font color to red
    title.text_frame.paragraphs[0].font.bold = True  # Set font to bold
    title.text_frame.paragraphs[0].font.italic = True  # Set font to bold

    # Left justify the text
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Make text left justified

    # Position and size of the image
    img_left = Inches(0.3)  # Horizontal position from the left
    img_top = Inches(0.7)    # Vertical position from the top
    img_height = Inches(4.8)  # Fixed height

    # Add the picture with automatic width based on aspect ratio
    image = slide.shapes.add_picture(img_path, img_left, img_top, height=img_height)
    
    # Calculate image width 
    # Convert to inches
    img_width_in = image.width.inches

    # Centralize the textbox in the remaining space
    txt_horiz = Inches(img_width_in + 0.3 + 0.15)
    txt_vert = Inches(0.31)
    txt_wid = Inches(4.76)
    txt_hei = Inches(4.54)

    # Add a textbox with the specified dimensions and position
    textbox = slide.shapes.add_textbox(txt_horiz, txt_vert, txt_wid, txt_hei)

    # TEXT frame
    text_frame = textbox.text_frame

    # Set vertical alignment to top
    # text_frame.margin_top = 0  # Ensure there's no top margin
    # text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Check if the length of novelas is less than 20
    # if len(novelas) <= 20:
    #     # Single column (helps not exceeding the width of the slide)
    #     # Fill col2 with blank entries
    #     col1 = novelas
    #     col2 = [''] * len(col1)  # Create a col2 with blank strings of the same length as col1
    # else:
    #     # Split the list into two equal parts
    #     split = len(novelas) // 2
    #     split = 25
    #     col1 = novelas[:split]  # First half for column 1
    #     col2 = novelas[split:]  # Second half for column 2
    #     # If col2 is shorter, fill it with blank strings
    #     # Ensure col2 matches the length of col1
    #     col2 = col2 + [''] * (len(col1) - len(col2))  
    
    #split = len(novelas) // 2
    split = 25
    col1 = novelas[:split]  # First half for column 1
    col2 = novelas[split:]  # Second half for column 2
    # If col2 is shorter, fill it with blank strings
    # Ensure col2 matches the length of col1
    col2 = col2 + [''] * (len(col1) - len(col2)) 
    
    # Determine the max length of names in col1 for consistent alignment
    max_len = max(len(name) for name in col1) + 1  # Adding extra space for padding

    # Add formatted text to the textbox
    for name1, name2 in zip(col1, col2):
        # Create a new paragraph for each entry
        paragraph = text_frame.add_paragraph()
        # TRY TO REMOVE BLANK LINE AT THE TOP
        # paragraph.space_before = 0
        if name2 != "":
           paragraph.text = f"• {name1.ljust(max_len)}• {name2}"
        else:
            paragraph.text = f"• {name1.ljust(max_len)}"   
        
        # Set font properties for the paragraph
        run = paragraph.runs[0]
        run.font.size = Pt(11)  # Set font size
        run.font.name = "Courier New"  # Set font type
        run.font.color.rgb = RGBColor(0, 0, 0)  # Set font color to black

        # If you want to make the text bold
        run.font.bold = True
    

# MAIN CODE
def Runs_cmd(PL_name=None,PL_nbr=None):
    # Load the Excel file into a DataFrame
    df = pd.read_excel("D:\\Videos\\Atrizes BR\\Atrizes.xlsx", sheet_name="Main")

    # OPENS EXCEL FILE
    Excel_file = "D:\\Videos\\Atrizes BR\\Atrizes.xlsx"
    excelf = Excel.open_excel(Excel_file,"Final")
    worksheet = excelf["sheet"]
    headers = excelf["headers"]

    rows = Excel.last_row(worksheet)-1

    # Artist	Title	Type
    Atriz_col = Excel.col_number(headers,"Atriz")
    selec_col = Excel.col_number(headers,"Select")
    Img_col = Excel.col_number(headers,"Image")
    
    next_row = 1
    slides = 0
    while next_row+1 <= rows+1: # and slides<=14
          next_row = next_row+1
          atriz = worksheet.cell(row=next_row, column=Atriz_col).value
          sel_value = worksheet.cell(row=next_row, column=selec_col).value
          Img_value = worksheet.cell(row=next_row, column=Img_col).value
          print("Creating slide",next_row-1,"of",rows,":", atriz)
          if sel_value=="x":
             novelas = read_novelas(df, atriz)
             add_slide(atriz, Img_value, novelas)
             slides = slides + 1
             #if next_row % 1==0:
                # Save the presentation
    presentation.save("D:\\Videos\\Atrizes BR\\New.pptx")

# CALLS FUNC
Runs_cmd()
