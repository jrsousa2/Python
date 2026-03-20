from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a slide with a title and content layout
slide_layout = presentation.slide_layouts[1]  # 1 is for 'Title and Content' layout
slide = presentation.slides.add_slide(slide_layout)

# Add title and content
title = slide.shapes.title
title.text = "My First Slide"

content = slide.placeholders[1]  # Content placeholder
content.text = "Hello, this is text added with Python!"

# Save the presentation
presentation.save("D:\\Videos\\Atrizes BR\\test.pptx")
